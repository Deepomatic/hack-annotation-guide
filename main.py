"""Generate a PowerPoint presentation outlining the plan for building
an Annotation Guide Generator skill/agent."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT_FILE = "annotation_guide_plan.pptx"

# ── Helpers ──────────────────────────────────────────────────────────


def _add_bullet_slide(prs, title_text, bullets, layout_index=1):
    """Add a slide with a title and a bulleted body."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    slide.shapes.title.text = title_text

    body = slide.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.space_after = Pt(6)
    return slide


def _add_two_column_slide(prs, title_text, left_items, right_items):
    """Add a slide with a title and two text-box columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    # Left column
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5))
    ltf = left.text_frame
    ltf.word_wrap = True
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(4)

    # Right column
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(5))
    rtf = right.text_frame
    rtf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(4)

    return slide


# ── Slide content ────────────────────────────────────────────────────


def build_presentation() -> Presentation:
    prs = Presentation()

    # 1 ── Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Annotation Guide Generator"
    slide.placeholders[1].text = (
        "A skill / agent that turns POC scoping documents & images\n"
        "into a ready-to-use annotation guide for our labelling team"
    )

    # 2 ── Problem statement
    _add_bullet_slide(prs, "Problem Statement", [
        "Every new annotation project requires a detailed guide for annotators.",
        "Guides are currently written by hand — slow, inconsistent, error-prone.",
        "POC scoping docs already contain class definitions, rules & edge cases.",
        "Sample images exist but are not systematically paired with instructions.",
        "→ We can automate this: docs + images → structured PPTX guide.",
    ])

    # 3 ── Inputs overview (two columns)
    _add_two_column_slide(
        prs,
        "Inputs",
        [
            "📄  POC Scoping Files",
            "─────────────────────",
            "• Class / label definitions",
            "• Annotation rules & guidelines",
            "• Edge-case descriptions",
            "• Acceptance criteria",
            "• Project metadata (name, version)",
        ],
        [
            "🖼️  Sample Images",
            "─────────────────────",
            "• Positive examples per class",
            "• Negative / edge-case examples",
            "• Bounding-box or segmentation previews",
            "• Before / after comparison pairs",
        ],
    )

    # 4 ── Architecture / Pipeline
    _add_bullet_slide(prs, "Pipeline Overview", [
        "Step 1 — Ingest & parse scoping documents (Markdown / PDF / DOCX)",
        "Step 2 — Extract annotation classes, rules and edge cases",
        "Step 3 — Match sample images to their corresponding classes",
        "Step 4 — (Optional) Use an LLM to enrich descriptions & rephrase rules",
        "Step 5 — Generate per-class slides: title, description, do/don't images",
        "Step 6 — Assemble final PPTX with table of contents & summary slide",
        "Step 7 — Export & distribute to annotation team",
    ])

    # 5 ── Output: the annotation guide
    _add_bullet_slide(prs, "Output: The Annotation Guide", [
        "One PPTX file ready to share with annotators.",
        "Cover slide with project name, version & date.",
        "Table of contents linking to each class section.",
        "Per-class slide(s):",
        "   • Class name & description",
        "   • ✅ Positive example images with captions",
        "   • ❌ Negative / edge-case images with captions",
        "   • Rules & tips for that class",
        "Summary / FAQ slide at the end.",
    ])

    # 6 ── Tech stack & tooling
    _add_bullet_slide(prs, "Tech Stack & Tooling", [
        "python-pptx — PowerPoint generation & templating",
        "LLM (GPT / Claude) — optional text enrichment & rephrasing",
        "Pillow — image resizing & thumbnail generation",
        "PyPDF2 / python-docx — scoping doc ingestion",
        "Click / Typer — CLI interface",
        "Pydantic — structured config & validation",
    ])

    # 7 ── Next steps
    _add_bullet_slide(prs, "Next Steps", [
        "1. Define a JSON / YAML schema for scoping-doc input",
        "2. Build the document parser (Markdown first, then PDF/DOCX)",
        "3. Create a slide template with branded styles",
        "4. Implement per-class slide generation with image placeholders",
        "5. Add optional LLM pass for description enrichment",
        "6. CLI wrapper: python main.py --scope spec.md --images ./imgs/",
        "7. Test on 2-3 real POC projects & iterate",
    ])

    return prs


# ── Main ─────────────────────────────────────────────────────────────


def main():
    prs = build_presentation()
    prs.save(OUTPUT_FILE)
    print(f"✅  Presentation saved to {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
