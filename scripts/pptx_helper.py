"""Reusable PPTX primitives and Studio helpers for annotation guide generation.

This module contains:
- Generic PPTX layout/style helpers (slides, text, shapes, images, grids)
- Image download helpers (sample images from Studio API)

All project-specific slide composition lives in build_pptx_slides.py.
Do NOT add project-specific slides here.
"""

import io
import logging
from pathlib import Path

from PIL import Image as PILImage

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Pt

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────────────────────────────
# Design System — Color Palette
# ──────────────────────────────────────────────────────────────────────

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_LIGHT = RGBColor(0x2D, 0x3E, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)
MUTED = RGBColor(0x88, 0x92, 0xA0)
DIVIDER = RGBColor(0xE1, 0xE5, 0xEB)
DARK_TEXT = RGBColor(0x1B, 0x2A, 0x4A)

# Accent colors per view kind
ORANGE = RGBColor(0xFF, 0x6B, 0x35)  # Detection
TEAL = RGBColor(0x2E, 0xC4, 0xB6)  # Classification
SKY_BLUE = RGBColor(0x5F, 0xA8, 0xD3)  # Tagging

# Good / Bad
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)

# Placeholder background
PLACEHOLDER_BG = RGBColor(0xF0, 0xF1, 0xF3)

# ──────────────────────────────────────────────────────────────────────
# Design System — Typography
# ──────────────────────────────────────────────────────────────────────

FONT_FAMILY = "Calibri"
FONT_FAMILY_MONO = "Consolas"

FONT_SIZE_COVER_TITLE = Pt(40)
FONT_SIZE_COVER_SUBTITLE = Pt(18)
FONT_SIZE_SECTION_TITLE = Pt(36)
FONT_SIZE_SECTION_SUBTITLE = Pt(16)
FONT_SIZE_SLIDE_TITLE = Pt(24)
FONT_SIZE_BODY = Pt(14)
FONT_SIZE_BODY_SMALL = Pt(12)
FONT_SIZE_CAPTION = Pt(10)
FONT_SIZE_LABEL = Pt(9)
FONT_SIZE_BADGE = Pt(11)
FONT_SIZE_FOOTER = Pt(8)

# ──────────────────────────────────────────────────────────────────────
# Design System — Layout (16:9 widescreen)
# ──────────────────────────────────────────────────────────────────────

SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

MARGIN_LEFT = Cm(2.0)
MARGIN_RIGHT = Cm(2.0)
MARGIN_TOP = Cm(1.5)
MARGIN_BOTTOM = Cm(1.2)

CONTENT_LEFT = MARGIN_LEFT
CONTENT_TOP = Cm(4.0)  # below title area
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM

TITLE_LEFT = MARGIN_LEFT
TITLE_TOP = Cm(1.2)
TITLE_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
TITLE_HEIGHT = Cm(2.0)

ACCENT_BAR_WIDTH = Cm(0.4)
ACCENT_BAR_HEIGHT = Cm(3.0)

# ──────────────────────────────────────────────────────────────────────
# Presentation factory
# ──────────────────────────────────────────────────────────────────────


def create_presentation() -> Presentation:
    """Create a blank 16:9 Presentation with no template."""
    prs = Presentation()
    prs.slide_width = int(SLIDE_WIDTH)
    prs.slide_height = int(SLIDE_HEIGHT)
    return prs


def add_blank_slide(prs: Presentation):
    """Add a blank slide (layout index 6 = Blank in default theme)."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


# ──────────────────────────────────────────────────────────────────────
# Background helpers
# ──────────────────────────────────────────────────────────────────────


def set_slide_bg_solid(slide, color: RGBColor):
    """Set a solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_slide_bg_gradient(slide, color_start: RGBColor, color_end: RGBColor):
    """Set a two-stop linear gradient background on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color_start
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color_end
    fill.gradient_stops[1].position = 1.0


# ──────────────────────────────────────────────────────────────────────
# Text helpers
# ──────────────────────────────────────────────────────────────────────


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size=FONT_SIZE_BODY,
    font_name=FONT_FAMILY,
    bold=False,
    italic=False,
    color: RGBColor = DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    word_wrap=True,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    """Add a text box with a single styled run. Returns the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    try:
        tf.vertical_anchor = vertical_anchor
    except Exception:
        pass

    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rich_textbox(
    slide,
    left,
    top,
    width,
    height,
    runs: list[dict],
    *,
    alignment=PP_ALIGN.LEFT,
    word_wrap=True,
    line_spacing=None,
):
    """Add a textbox with multiple styled runs in a single paragraph.

    Each run dict: {"text": str, "bold": bool, "color": RGBColor, "size": Pt, "font": str}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = alignment
    if line_spacing is not None:
        para.line_spacing = line_spacing

    for r in runs:
        run = para.add_run()
        run.text = r.get("text", "")
        run.font.size = r.get("size", FONT_SIZE_BODY)
        run.font.name = r.get("font", FONT_FAMILY)
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", DARK_TEXT)
    return txBox


def add_multiline_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[dict],
    *,
    word_wrap=True,
    para_spacing_pt=4,
):
    """Add a textbox with multiple paragraphs (one per line).

    Each line dict: {"text": str, "bold": bool, "color": RGBColor,
                     "size": Pt, "alignment": PP_ALIGN, "indent_level": int}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = line.get("alignment", PP_ALIGN.LEFT)
        para.space_after = Pt(para_spacing_pt)
        run = para.add_run()
        run.text = line.get("text", "")
        run.font.size = line.get("size", FONT_SIZE_BODY)
        run.font.name = line.get("font", FONT_FAMILY)
        run.font.bold = line.get("bold", False)
        run.font.italic = line.get("italic", False)
        run.font.color.rgb = line.get("color", DARK_TEXT)
    return txBox


def add_title(slide, text: str, *, color: RGBColor = DARK_TEXT, font_size=FONT_SIZE_SLIDE_TITLE):
    """Add a title textbox at the standard title position."""
    return add_textbox(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        text,
        font_size=font_size,
        bold=True,
        color=color,
    )


def add_subtitle(slide, text: str, *, color: RGBColor = MUTED, font_size=FONT_SIZE_BODY):
    """Add subtitle text just below the standard title position."""
    return add_textbox(
        slide,
        TITLE_LEFT,
        TITLE_TOP + Cm(2.2),
        TITLE_WIDTH,
        Cm(1.2),
        text,
        font_size=font_size,
        color=color,
    )


# ──────────────────────────────────────────────────────────────────────
# Shape helpers
# ──────────────────────────────────────────────────────────────────────


def add_rounded_rect(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color: RGBColor | None = WHITE,
    border_color: RGBColor | None = DIVIDER,
    border_width=Pt(1.5),
    text: str = "",
    text_color: RGBColor = DARK_TEXT,
    text_size=FONT_SIZE_BODY_SMALL,
    text_bold=False,
    text_align=PP_ALIGN.CENTER,
):
    """Add a rounded rectangle shape with optional text."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if border_color is not None:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = text_align
        run = para.add_run()
        run.text = text
        run.font.size = text_size
        run.font.color.rgb = text_color
        run.font.bold = text_bold
        run.font.name = FONT_FAMILY

    return shape


def add_accent_bar(slide, left, top, width, height, color: RGBColor):
    """Add a solid colored rectangle (accent bar / decorative element)."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, *, color: RGBColor = DIVIDER, width=Pt(1)):
    """Add a straight connector line."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        x1, y1, x2, y2,
    )
    conn.line.color.rgb = color
    conn.line.width = width
    return conn


def add_badge(
    slide,
    left,
    top,
    text: str,
    *,
    bg_color: RGBColor,
    text_color: RGBColor = WHITE,
    font_size=FONT_SIZE_BADGE,
    padding_lr=Cm(0.6),
    height=Cm(1.0),
):
    """Add a pill-shaped badge (rounded rect with text)."""
    # Estimate width from text length
    char_width = int(font_size) * 0.7
    text_width = int(len(text) * char_width)
    width = max(text_width + 2 * int(padding_lr), int(Cm(2.5)))

    return add_rounded_rect(
        slide,
        left, top,
        width, height,
        fill_color=bg_color,
        border_color=None,
        text=text,
        text_color=text_color,
        text_size=font_size,
        text_bold=True,
    )


# ──────────────────────────────────────────────────────────────────────
# Image helpers
# ──────────────────────────────────────────────────────────────────────


def add_image(slide, img_path, left, top, max_width, max_height):
    """Add an image auto-scaled to fit within max_width × max_height.

    Preserves aspect ratio. Centers within the given area.
    Returns the picture shape, or None if image can't be loaded.
    """
    from PIL import Image as PILImage

    img_path = Path(img_path)
    if not img_path.exists():
        return None

    try:
        with PILImage.open(img_path) as im:
            iw, ih = im.size
    except Exception:
        return None

    if iw <= 0 or ih <= 0:
        return None

    scale = min(int(max_width) / iw, int(max_height) / ih)
    pic_w = int(iw * scale)
    pic_h = int(ih * scale)

    # Center in area
    x_off = (int(max_width) - pic_w) // 2
    y_off = (int(max_height) - pic_h) // 2

    try:
        pic = slide.shapes.add_picture(
            str(img_path),
            left + x_off,
            top + y_off,
            pic_w,
            pic_h,
        )
        return pic
    except Exception as exc:
        logger.warning("Could not embed %s: %s", img_path.name, exc)
        return None


def add_image_placeholder(
    slide,
    left,
    top,
    width,
    height,
    label: str = "",
    *,
    border_color: RGBColor = DIVIDER,
    bg_color: RGBColor = PLACEHOLDER_BG,
    label_color: RGBColor | None = None,
    dash=True,
):
    """Draw a labeled rectangle as an image placeholder."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    if dash:
        shape.line.dash_style = 2  # dashed

    if label:
        tf = shape.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label
        run.font.size = FONT_SIZE_CAPTION
        run.font.color.rgb = label_color or border_color
        run.font.bold = True
        run.font.name = FONT_FAMILY

    return shape


# ──────────────────────────────────────────────────────────────────────
# Layout helpers
# ──────────────────────────────────────────────────────────────────────


def grid_positions(
    n_items: int,
    area_left,
    area_top,
    area_width,
    area_height,
    *,
    cols: int | None = None,
    h_padding=Cm(0.5),
    v_padding=Cm(0.5),
) -> list[tuple[int, int, int, int]]:
    """Compute (left, top, width, height) for N items in a grid layout.

    Returns list of (left, top, cell_width, cell_height).
    If cols is None, auto-picks: 1 for 1 item, 2 for 2, 3 for 3-6, 4 for 7+.
    """
    if n_items == 0:
        return []

    if cols is None:
        if n_items <= 1:
            cols = 1
        elif n_items <= 2:
            cols = 2
        elif n_items <= 6:
            cols = 3
        else:
            cols = 4

    import math
    rows = math.ceil(n_items / cols)

    cell_w = (int(area_width) - (cols - 1) * int(h_padding)) // cols
    cell_h = (int(area_height) - (rows - 1) * int(v_padding)) // rows

    positions = []
    for i in range(n_items):
        row = i // cols
        col = i % cols
        x = int(area_left) + col * (cell_w + int(h_padding))
        y = int(area_top) + row * (cell_h + int(v_padding))
        positions.append((x, y, cell_w, cell_h))

    return positions


# ──────────────────────────────────────────────────────────────────────
# Connector helpers (for tree/graph diagrams)
# ──────────────────────────────────────────────────────────────────────


def add_tree_connector(slide, parent_cx, parent_bottom_y, child_cx, child_top_y, *, color=DIVIDER, width=Pt(1.5)):
    """Draw a connector from parent bottom-center to child top-center."""
    return add_line(slide, parent_cx, parent_bottom_y, child_cx, child_top_y, color=color, width=width)


# ──────────────────────────────────────────────────────────────────────
# Footer / page number
# ──────────────────────────────────────────────────────────────────────


def add_footer(slide, text: str, *, include_page_number=False, page_number: int | None = None):
    """Add a small footer text at the bottom-right of the slide."""
    footer_text = text
    if include_page_number and page_number is not None:
        footer_text = f"{text}  |  {page_number}"

    return add_textbox(
        slide,
        SLIDE_WIDTH - Cm(8),
        SLIDE_HEIGHT - Cm(1.0),
        Cm(7),
        Cm(0.8),
        footer_text,
        font_size=FONT_SIZE_FOOTER,
        color=MUTED,
        alignment=PP_ALIGN.RIGHT,
    )


# ──────────────────────────────────────────────────────────────────────
# Decorative helpers
# ──────────────────────────────────────────────────────────────────────


def add_bottom_strip(slide, color: RGBColor, height=Cm(0.3)):
    """Add a thin colored strip at the bottom of a slide."""
    return add_accent_bar(
        slide,
        0, SLIDE_HEIGHT - height,
        SLIDE_WIDTH, height,
        color,
    )


def add_top_line(slide, *, color: RGBColor = DIVIDER, y=None):
    """Add a thin horizontal line below the title area (visual separator)."""
    if y is None:
        y = TITLE_TOP + TITLE_HEIGHT + Cm(0.3)
    return add_line(
        slide,
        MARGIN_LEFT, y,
        SLIDE_WIDTH - MARGIN_RIGHT, y,
        color=color,
        width=Pt(0.75),
    )


# ──────────────────────────────────────────────────────────────────────
# Card / panel helpers (modern look)
# ──────────────────────────────────────────────────────────────────────


def add_card(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color: RGBColor = WHITE,
    border_color: RGBColor | None = DIVIDER,
    border_width=Pt(0.75),
    shadow=False,
):
    """Add a card panel (rounded rect with subtle border). Returns the shape."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    if shadow:
        # Add a simple shadow effect via XML
        try:
            sp = shape._element
            spPr = sp.find(qn("a:spPr"))
            if spPr is None:
                from lxml import etree
                spPr = sp.find(qn("p:spPr"))
            if spPr is not None:
                effectLst = spPr.find(qn("a:effectLst"))
                if effectLst is None:
                    from lxml import etree
                    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
                outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
                outerShdw.set("blurRad", "50800")
                outerShdw.set("dist", "25400")
                outerShdw.set("dir", "5400000")
                outerShdw.set("algn", "tl")
                srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
                srgbClr.set("val", "000000")
                alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha.set("val", "15000")
        except Exception:
            pass  # shadow is cosmetic, don't fail

    return shape


def add_vertical_divider(slide, x, top, height, *, color: RGBColor = DIVIDER, width=Pt(0.75)):
    """Add a vertical divider line."""
    return add_line(slide, x, top, x, top + height, color=color, width=width)


# ──────────────────────────────────────────────────────────────────────
# Project map parsing utilities
# ──────────────────────────────────────────────────────────────────────

# View-kind labels & accent colors
KIND_LABELS = {
    "CLA": "Classification",
    "DET": "Detection",
    "TAG": "Tagging",
}


def kind_color(kind: str) -> RGBColor:
    """Return the accent color for a view kind."""
    if kind == "DET":
        return ORANGE
    if kind == "CLA":
        return TEAL
    return SKY_BLUE


def sanitize_name(name: str) -> str:
    """Turn a label into a safe directory/filename component."""
    return name.replace(" ", "_").replace("/", "-").replace("\\", "-")


def build_tree(project_map: dict):
    """Parse project map into (nodes_dict, roots_list).

    Each node: {id, label, kind, parent, conditions, tag_names, children}
    """
    nodes = {}
    for n in project_map["nodes"]:
        nodes[n["id"]] = {
            "id": n["id"],
            "label": n["label"],
            "kind": n["data"].get("kind", ""),
            "parent": n["data"].get("parent", ""),
            "conditions": n["data"].get("conditions") or [],
            "tag_names": n["data"].get("tag_names", []),
            "children": [],
        }

    for edge in project_map["edges"]:
        src, tgt = edge["source"], edge["target"]
        if src and src in nodes and tgt in nodes:
            nodes[src]["children"].append(tgt)

    roots = [nid for nid, n in nodes.items() if not n["parent"]]
    return nodes, roots


def build_concept_map(project_map: dict) -> dict:
    """Return dict {concept_id: concept_name}."""
    return {c["id"]: c["concept_name"] for c in project_map.get("concepts", [])}


def resolve_conditions(conditions, concept_map: dict) -> str:
    """Turn condition lists into a human-readable string."""
    if not conditions:
        return ""
    groups = []
    for group in conditions:
        names = [concept_map.get(cid, str(cid)) for cid in group]
        groups.append(" & ".join(names))
    return "  |  ".join(groups)


def dfs_order(nodes: dict, roots: list) -> list[str]:
    """Return node IDs in depth-first order."""
    order = []
    visited = set()

    def dfs(nid):
        if nid in visited:
            return
        visited.add(nid)
        order.append(nid)
        for child_id in nodes[nid]["children"]:
            dfs(child_id)

    for r in roots:
        dfs(r)
    for nid in nodes:
        if nid not in visited:
            order.append(nid)
    return order


def find_view_images(view_label: str, images_dir) -> dict[str, list[Path]]:
    """Return {concept_name: [path1, path2, ...]} for a view's downloaded images.

    Files expected as: {view}__{concept}__{idx}.ext  or  {view}__{concept}.ext
    """
    if images_dir is None:
        return {}
    view_dir = Path(images_dir) / sanitize_name(view_label)
    if not view_dir.is_dir():
        return {}
    exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}
    result: dict[str, list[Path]] = {}
    for p in sorted(view_dir.iterdir()):
        if p.suffix.lower() not in exts:
            continue
        stem = p.stem
        if "__" in stem:
            parts = stem.split("__")
            concept = parts[1].replace("_", " ") if len(parts) >= 2 else stem
        else:
            concept = stem
        result.setdefault(concept, []).append(p)
    return result


def match_images(concept_name: str, view_images: dict[str, list[Path]]) -> list[Path]:
    """Match a concept name to downloaded images. Returns list of paths."""
    if not view_images:
        return []
    if concept_name in view_images:
        return view_images[concept_name]
    sanitized = sanitize_name(concept_name).replace("_", " ")
    for k, v in view_images.items():
        if k == sanitized or k.lower() == concept_name.lower():
            return v
    cn_lower = concept_name.lower().replace(" ", "_")
    for k, v in view_images.items():
        if cn_lower in k.lower().replace(" ", "_"):
            return v
    return []


# ──────────────────────────────────────────────────────────────────────
# Tree layout computation (for overview diagram)
# ──────────────────────────────────────────────────────────────────────


def compute_tree_positions(nodes, roots):
    """Return {nid: (grid_x, grid_y)}, total_width, num_levels."""
    widths: dict[str, float] = {}

    def calc_width(nid):
        children = [c for c in nodes[nid]["children"] if c in nodes]
        if not children:
            widths[nid] = 1.0
        else:
            widths[nid] = sum(calc_width(c) for c in children)
        return widths[nid]

    total_w = sum(calc_width(r) for r in roots)

    positions: dict[str, tuple[float, int]] = {}
    max_depth = 0

    def assign(nid, x_start, depth):
        nonlocal max_depth
        max_depth = max(max_depth, depth)
        w = widths[nid]
        positions[nid] = (x_start + w / 2.0, depth)
        children = [c for c in nodes[nid]["children"] if c in nodes]
        cx = x_start
        for c in children:
            assign(c, cx, depth + 1)
            cx += widths[c]

    x = 0.0
    for r in roots:
        assign(r, x, 0)
        x += widths[r]

    return positions, total_w, max_depth + 1


# ══════════════════════════════════════════════════════════════════════
# Pre-built slide templates
# ══════════════════════════════════════════════════════════════════════
#
# These are ready-to-use, configurable slide builders.
# Import and call them from build_pptx_slides.py (or any script).
#
# USAGE EXAMPLE — build_concept_recap_slide
# ──────────────────────────────────────────
#
#   from pptx_helper import (
#       create_presentation, build_concept_recap_slide,
#       build_tree, find_view_images, ORANGE,
#   )
#
#   prs = create_presentation()
#   nodes, roots = build_tree(project_map)
#   node = nodes[some_view_id]
#
#   # Default style:
#   build_concept_recap_slide(prs, node, images_dir="images/myproject")
#
#   # Custom title, accent, background:
#   build_concept_recap_slide(
#       prs, node,
#       images_dir="images/myproject",
#       title="LED Indicators — Overview",
#       accent_color=ORANGE,
#       bg_color=RGBColor(0xFF, 0xFF, 0xFF),
#       max_concepts=6,
#   )
#
#   # Same pattern for all other build_* functions — each accepts
#   # keyword overrides for title, colors, and layout knobs.
#
# ══════════════════════════════════════════════════════════════════════


def build_cover_slide(
    prs,
    project_name: str = "",
    *,
    title: str = "Annotation Guide",
    bg_start: RGBColor = None,
    bg_end: RGBColor = None,
    accent_color: RGBColor = None,
    title_color: RGBColor = None,
    subtitle_color: RGBColor = None,
):
    """Full-bleed cover slide with gradient background.

    Args:
        title: Main title text.
        project_name: Subtitle (project/org name).
        bg_start/bg_end: Gradient colors (default NAVY → NAVY_LIGHT).
        accent_color: Top/bottom strip and side bar color.
        title_color/subtitle_color: Override text colors.
    """
    bg_start = bg_start or NAVY
    bg_end = bg_end or NAVY_LIGHT
    accent = accent_color or ORANGE
    tc = title_color or WHITE
    sc = subtitle_color or MUTED

    slide = add_blank_slide(prs)
    set_slide_bg_gradient(slide, bg_start, bg_end)

    add_accent_bar(slide, 0, 0, SLIDE_WIDTH, Cm(0.2), accent)
    add_accent_bar(slide, Cm(3.0), Cm(5.0), Cm(0.3), Cm(4.5), accent)

    add_textbox(slide, Cm(4.5), Cm(5.5), Cm(25), Cm(3.0), title,
                font_size=FONT_SIZE_COVER_TITLE, bold=True, color=tc)

    if project_name:
        add_textbox(slide, Cm(4.5), Cm(9.0), Cm(25), Cm(2.0), project_name,
                    font_size=FONT_SIZE_COVER_SUBTITLE, color=sc)

    add_line(slide, Cm(4.5), Cm(12.0), Cm(15), Cm(12.0), color=MUTED, width=Pt(0.5))
    add_accent_bar(slide, 0, SLIDE_HEIGHT - Cm(0.2), SLIDE_WIDTH, Cm(0.2), accent)
    return slide


def build_toc_slide(
    prs,
    nodes: dict,
    roots: list,
    *,
    title: str = "Table of Contents",
    bg_color: RGBColor = None,
    header_line_color: RGBColor = None,
    strip_color: RGBColor = None,
):
    """Table of Contents slide listing root views with kind badges.

    Args:
        nodes/roots: Parsed project tree (from build_tree).
        title: Slide title.
        bg_color: Slide background.
        header_line_color: Color of line under title.
        strip_color: Bottom strip color.
    """
    bg = bg_color or LIGHT_BG
    hline = header_line_color or NAVY
    strip = strip_color or NAVY

    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, bg)
    add_title(slide, title)
    add_top_line(slide, color=hline)

    y = CONTENT_TOP + Cm(0.5)
    for root_id in roots:
        node = nodes[root_id]
        accent = kind_color(node["kind"])
        kind_lbl = KIND_LABELS.get(node["kind"], node["kind"])
        n_concepts = len(node.get("tag_names", []))

        add_card(slide, MARGIN_LEFT, y - Cm(0.1), CONTENT_WIDTH, Cm(1.3),
                 fill_color=WHITE, border_color=DIVIDER, shadow=True)
        add_accent_bar(slide, MARGIN_LEFT, y - Cm(0.1), Cm(0.25), Cm(1.3), accent)
        add_textbox(slide, MARGIN_LEFT + Cm(1.0), y, Cm(18), Cm(0.8),
                    node["label"], font_size=FONT_SIZE_BODY, bold=True, color=DARK_TEXT)
        add_badge(slide, MARGIN_LEFT + Cm(20), y + Cm(0.05), kind_lbl,
                  bg_color=accent, text_color=WHITE, height=Cm(0.8), font_size=FONT_SIZE_LABEL)
        if n_concepts:
            add_textbox(slide, SLIDE_WIDTH - MARGIN_RIGHT - Cm(5), y, Cm(4.5), Cm(0.8),
                        f"{n_concepts} concepts", font_size=FONT_SIZE_CAPTION,
                        color=MUTED, alignment=PP_ALIGN.RIGHT)
        y += Cm(1.6)

    add_bottom_strip(slide, strip)
    return slide


def build_overview_slide(
    prs,
    nodes: dict,
    roots: list,
    concept_map: dict,
    *,
    title: str = "Views Overview",
    bg_color: RGBColor = None,
    header_line_color: RGBColor = None,
    connector_color: RGBColor = None,
    strip_color: RGBColor = None,
):
    """Tree diagram overview of the full view hierarchy.

    Args:
        nodes/roots/concept_map: From build_tree / build_concept_map.
        title: Slide title.
        connector_color: Color of tree edges.
    """
    bg = bg_color or LIGHT_BG
    hline = header_line_color or NAVY
    conn_c = connector_color or MUTED
    strip = strip_color or NAVY

    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, bg)
    add_title(slide, title)
    add_top_line(slide, color=hline)

    positions, total_w, num_levels = compute_tree_positions(nodes, roots)
    if not positions:
        return slide

    area_left = Cm(1.5)
    area_top = Cm(4.5)
    area_w = SLIDE_WIDTH - Cm(3.0)
    area_h = SLIDE_HEIGHT - area_top - Cm(2.0)

    node_w = min(Cm(4.5), int(int(area_w) / max(total_w, 1) - int(Cm(0.6))))
    node_w = max(node_w, Cm(2.2))
    node_h = Cm(1.4)

    cell_w = int(area_w) / total_w if total_w else int(area_w)
    cell_h = int(area_h) / num_levels if num_levels > 1 else int(area_h)

    emu_pos: dict[str, tuple[int, int]] = {}
    for nid, (gx, gy) in positions.items():
        cx = int(int(area_left) + gx * cell_w)
        cy = int(int(area_top) + gy * cell_h + int(node_h) // 2)
        emu_pos[nid] = (cx, cy)

    # Connectors
    for nid in positions:
        children = [c for c in nodes[nid]["children"] if c in nodes]
        if not children:
            continue
        px, py = emu_pos[nid]
        parent_bottom_y = py + int(node_h) // 2
        for child_id in children:
            ccx, ccy = emu_pos[child_id]
            child_top_y = ccy - int(node_h) // 2
            add_tree_connector(slide, px, parent_bottom_y, ccx, child_top_y,
                               color=conn_c, width=Pt(1.5))
            cond_text = resolve_conditions(nodes[child_id]["conditions"], concept_map)
            if cond_text:
                mid_x = (px + ccx) // 2
                mid_y = (parent_bottom_y + child_top_y) // 2
                lbl_w = min(Cm(5.0), int(cell_w * 0.95))
                add_textbox(slide, mid_x - lbl_w // 2, mid_y - Cm(0.4), lbl_w, Cm(0.8),
                            f"({cond_text})", font_size=Pt(7), color=MUTED,
                            alignment=PP_ALIGN.CENTER)

    # Node boxes
    for nid, (cx, cy) in emu_pos.items():
        n = nodes[nid]
        accent = kind_color(n["kind"])
        add_card(slide, cx - int(node_w) // 2, cy - int(node_h) // 2, node_w, node_h,
                 fill_color=WHITE, border_color=accent, border_width=Pt(2), shadow=True)
        add_accent_bar(slide, cx - int(node_w) // 2, cy - int(node_h) // 2,
                       node_w, Cm(0.2), accent)
        add_textbox(slide, cx - int(node_w) // 2, cy - int(node_h) // 2 + Cm(0.25),
                    node_w, int(node_h) - Cm(0.25), n["label"],
                    font_size=Pt(9), bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Legend
    legend_y = SLIDE_HEIGHT - Cm(1.6)
    legend_x = MARGIN_LEFT
    for kind_code, kind_label in KIND_LABELS.items():
        c = kind_color(kind_code)
        add_accent_bar(slide, legend_x, legend_y + Cm(0.1), Cm(0.6), Cm(0.15), c)
        add_textbox(slide, legend_x + Cm(0.8), legend_y, Cm(4), Cm(0.6),
                    kind_label, font_size=FONT_SIZE_LABEL, color=MUTED)
        legend_x += Cm(5)

    add_bottom_strip(slide, strip)
    return slide


def build_section_slide(
    prs,
    view_name: str,
    view_kind: str,
    *,
    accent_color: RGBColor = None,
    bg_start: RGBColor = None,
    bg_end: RGBColor = None,
    title_color: RGBColor = None,
):
    """Section divider slide — navy gradient, left accent bar, kind badge.

    Args:
        view_name: Large title text.
        view_kind: One of "DET", "CLA", "TAG".
        accent_color: Override per-kind accent.
        bg_start/bg_end: Gradient background.
    """
    bg_s = bg_start or NAVY
    bg_e = bg_end or NAVY_LIGHT
    accent = accent_color or kind_color(view_kind)
    tc = title_color or WHITE

    slide = add_blank_slide(prs)
    set_slide_bg_gradient(slide, bg_s, bg_e)

    add_accent_bar(slide, 0, Cm(4.5), Cm(0.35), Cm(10), accent)
    add_line(slide, Cm(3.5), Cm(11.5), Cm(20), Cm(11.5), color=MUTED, width=Pt(0.5))
    add_textbox(slide, Cm(3.5), Cm(5.5), Cm(26), Cm(3.5), view_name,
                font_size=FONT_SIZE_SECTION_TITLE, bold=True, color=tc)

    kind_label = KIND_LABELS.get(view_kind, view_kind)
    add_badge(slide, Cm(3.5), Cm(10.0), kind_label, bg_color=accent, text_color=WHITE)
    add_accent_bar(slide, 0, SLIDE_HEIGHT - Cm(0.2), SLIDE_WIDTH, Cm(0.2), accent)
    return slide


def build_info_slide(
    prs,
    node: dict,
    nodes: dict,
    concept_map: dict,
    *,
    bg_color: RGBColor = None,
    accent_color: RGBColor = None,
    instruction_text: str = "Add annotation instructions for this view here.",
    extra_fields: list[tuple[str, str]] | None = None,
):
    """Info slide with key-value metadata in a card.

    Args:
        node/nodes/concept_map: View data.
        accent_color: Override per-kind accent.
        instruction_text: Placeholder text below the info card.
        extra_fields: Additional (label, value) pairs appended after defaults.
    """
    bg = bg_color or LIGHT_BG
    kind = node["kind"]
    accent = accent_color or kind_color(kind)

    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, bg)

    add_title(slide, node["label"])
    kind_label = KIND_LABELS.get(kind, kind)
    add_badge(slide, SLIDE_WIDTH - MARGIN_RIGHT - Cm(5), TITLE_TOP + Cm(0.3),
              kind_label, bg_color=accent, text_color=WHITE)
    add_top_line(slide, color=accent)

    card_left, card_top = MARGIN_LEFT, CONTENT_TOP
    card_w, card_h = CONTENT_WIDTH, Cm(7.5)

    add_card(slide, card_left, card_top, card_w, card_h,
             fill_color=WHITE, border_color=DIVIDER, shadow=True)
    add_accent_bar(slide, card_left, card_top, ACCENT_BAR_WIDTH, card_h, accent)

    info_left = card_left + Cm(1.5)
    info_w = card_w - Cm(2.5)
    y = card_top + Cm(0.6)
    row_h = Cm(1.4)

    parent_label = ""
    if node["parent"] and node["parent"] in nodes:
        parent_label = nodes[node["parent"]]["label"]
    cond_str = resolve_conditions(node["conditions"], concept_map)
    child_labels = [nodes[c]["label"] for c in node["children"] if c in nodes]
    tag_names = node.get("tag_names", [])

    fields: list[tuple[str, str]] = [
        ("Parent view", parent_label or "— (root)"),
        ("Activated by", cond_str or "— (always)"),
        ("Child views", ", ".join(child_labels) if child_labels else "— (none)"),
    ]
    if tag_names:
        fields.append(("Concepts", ", ".join(tag_names)))
    if extra_fields:
        fields.extend(extra_fields)

    for i, (label, value) in enumerate(fields):
        add_textbox(slide, info_left, y, Cm(6), Cm(0.6), label.upper(),
                    font_size=FONT_SIZE_CAPTION, bold=True, color=MUTED)
        add_textbox(slide, info_left + Cm(6.5), y, info_w - Cm(6.5), Cm(0.6),
                    value, font_size=FONT_SIZE_BODY, color=DARK_TEXT)
        if i < len(fields) - 1:
            add_line(slide, info_left, y + row_h - Cm(0.2),
                     info_left + info_w, y + row_h - Cm(0.2),
                     color=DIVIDER, width=Pt(0.5))
        y += row_h

    if instruction_text:
        add_textbox(slide, card_left + Cm(1.5), card_top + card_h + Cm(0.8),
                    card_w - Cm(3), Cm(3), instruction_text,
                    font_size=FONT_SIZE_BODY, italic=True, color=MUTED)

    add_bottom_strip(slide, accent)
    return slide


def build_concept_recap_slide(
    prs,
    node: dict,
    images_dir=None,
    *,
    title: str | None = None,
    accent_color: RGBColor = None,
    bg_color: RGBColor = None,
    max_concepts: int = 9,
    concepts_override: list[str] | None = None,
):
    """Recap slide showing all concepts for a view in a grid.

    Layout: 1–3 → row, 4–6 → 3×2, 7–9 → 3×3, >max → text list.

    Args:
        title: Override slide title (default: "{view} — Concepts Overview").
        accent_color: Override per-kind accent.
        bg_color: Slide background.
        max_concepts: Cap before falling back to text list (default 9).
        concepts_override: Explicit list of concept names (skips tag_names).
    """
    bg = bg_color or LIGHT_BG
    kind = node["kind"]
    accent = accent_color or kind_color(kind)
    slide_title = title or f"{node['label']}  —  Concepts Overview"

    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, bg)
    add_title(slide, slide_title)
    add_top_line(slide, color=accent)

    view_images = find_view_images(node["label"], images_dir)
    if concepts_override is not None:
        concepts = list(concepts_override)
    else:
        tag_names = node.get("tag_names", [])
        concepts = list(tag_names) if tag_names else list(view_images.keys()) if view_images else []

    if not concepts:
        add_textbox(slide, MARGIN_LEFT, CONTENT_TOP + Cm(3), CONTENT_WIDTH, Cm(2),
                    "No concepts defined for this view.",
                    font_size=FONT_SIZE_BODY, italic=True, color=MUTED,
                    alignment=PP_ALIGN.CENTER)
        add_bottom_strip(slide, accent)
        return slide

    if len(concepts) > max_concepts:
        add_textbox(slide, MARGIN_LEFT, CONTENT_TOP + Cm(2), CONTENT_WIDTH, Cm(2),
                    f"This view has {len(concepts)} concepts — too many for a single overview.",
                    font_size=FONT_SIZE_BODY, color=MUTED, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, MARGIN_LEFT + Cm(2), CONTENT_TOP + Cm(4.5),
                    CONTENT_WIDTH - Cm(4), Cm(6), ", ".join(concepts),
                    font_size=FONT_SIZE_BODY_SMALL, color=DARK_TEXT)
        add_bottom_strip(slide, accent)
        return slide

    n = len(concepts)
    cols = n if n <= 3 else 3

    gallery_top = CONTENT_TOP + Cm(0.3)
    gallery_h = SLIDE_HEIGHT - gallery_top - Cm(1.5)

    positions = grid_positions(n, MARGIN_LEFT, gallery_top, CONTENT_WIDTH, gallery_h,
                               cols=cols, h_padding=Cm(0.6), v_padding=Cm(0.5))

    for (x, y, w, h), concept_name in zip(positions, concepts):
        label_h = Cm(1.1)
        card_h = h
        img_h = card_h - label_h - Cm(0.6)
        pad = Cm(0.3)

        add_card(slide, x, y, w, card_h, fill_color=WHITE, border_color=DIVIDER, shadow=True)
        add_accent_bar(slide, x, y, w, Cm(0.15), accent)

        img_paths = match_images(concept_name, view_images)
        ix, iy, iw, ih = x + pad, y + Cm(0.3), w - 2 * pad, img_h

        if img_paths and img_paths[0].exists():
            add_image(slide, img_paths[0], ix, iy, iw, ih)
        else:
            add_image_placeholder(slide, ix, iy, iw, ih, label="[ sample ]",
                                  border_color=DIVIDER, bg_color=PLACEHOLDER_BG)

        add_textbox(slide, x, y + card_h - label_h, w, label_h, concept_name,
                    font_size=FONT_SIZE_BODY_SMALL, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)

    add_bottom_strip(slide, accent)
    return slide


def build_concept_detail_slide(
    prs,
    node: dict,
    concept_name: str,
    images_dir=None,
    *,
    accent_color: RGBColor = None,
    bg_color: RGBColor = None,
    good_color: RGBColor = None,
    bad_color: RGBColor = None,
    good_label: str = "✓  Good Examples",
    bad_label: str = "✗  Bad Examples",
    good_images: list | None = None,
    bad_images: list | None = None,
    n_slots: int = 2,
):
    """Per-concept split slide: good examples (left) / bad examples (right).

    Args:
        concept_name: Concept being illustrated.
        accent_color: Top-line accent (default per-kind).
        good_color/bad_color: Left/right panel stripes.
        good_label/bad_label: Section headers.
        good_images: Explicit list of Path for good side (overrides auto-match).
        bad_images: Explicit list of Path for bad side.
        n_slots: Number of image slots per side (stacked vertically).
    """
    bg = bg_color or LIGHT_BG
    accent = accent_color or kind_color(node["kind"])
    gc = good_color or GREEN
    bc = bad_color or RED

    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, bg)

    add_textbox(slide, TITLE_LEFT, TITLE_TOP, Cm(25), TITLE_HEIGHT, concept_name,
                font_size=FONT_SIZE_SLIDE_TITLE, bold=True, color=DARK_TEXT)
    add_textbox(slide, SLIDE_WIDTH - MARGIN_RIGHT - Cm(10), TITLE_TOP + Cm(0.2),
                Cm(9.5), Cm(1.0), node["label"],
                font_size=FONT_SIZE_CAPTION, color=MUTED, alignment=PP_ALIGN.RIGHT)
    add_top_line(slide, color=accent)

    half_w = (int(CONTENT_WIDTH) - int(Cm(1.0))) // 2
    left_x = int(MARGIN_LEFT)
    right_x = int(MARGIN_LEFT) + half_w + int(Cm(1.0))
    area_top = int(CONTENT_TOP) + int(Cm(0.3))
    area_h = int(SLIDE_HEIGHT) - area_top - int(Cm(1.5))
    pad = int(Cm(0.5))

    # Resolve images
    if good_images is None:
        view_images = find_view_images(node["label"], images_dir)
        good_images = match_images(concept_name, view_images)
    if bad_images is None:
        bad_images = []

    def _draw_panel(panel_x, panel_color, label_text, images):
        add_card(slide, panel_x, area_top, half_w, area_h,
                 fill_color=WHITE, border_color=DIVIDER, shadow=True)
        add_accent_bar(slide, panel_x, area_top, half_w, Cm(0.15), panel_color)
        add_textbox(slide, panel_x + Cm(0.5), area_top + Cm(0.3),
                    half_w - Cm(1), Cm(1.0), label_text,
                    font_size=FONT_SIZE_BODY, bold=True, color=panel_color)

        slot_top = area_top + Cm(1.6)
        slot_gap = int(Cm(0.3))
        slot_h = (area_h - Cm(2.0) - slot_gap * (n_slots - 1)) // n_slots
        slot_w = half_w - 2 * pad

        for i in range(n_slots):
            sx = panel_x + pad
            sy = slot_top + i * (int(slot_h) + slot_gap)

            if i < len(images) and images[i].exists():
                add_image(slide, images[i], sx, sy, slot_w, slot_h)
            else:
                placeholder_lbl = f"[ {label_text.split()[-1].lower()} ]"
                add_image_placeholder(slide, sx, sy, slot_w, slot_h,
                                      label=placeholder_lbl, border_color=panel_color,
                                      bg_color=PLACEHOLDER_BG)

    _draw_panel(left_x, gc, good_label, good_images)

    divider_x = left_x + half_w + int(Cm(0.5))
    add_vertical_divider(slide, divider_x, area_top, area_h, color=DIVIDER)

    _draw_panel(right_x, bc, bad_label, bad_images)

    add_bottom_strip(slide, accent)
    return slide


# ══════════════════════════════════════════════════════════════════════
# Image download helpers (from Studio API)
# ══════════════════════════════════════════════════════════════════════


def _dl_sanitize(name: str) -> str:
    """Turn a label into a safe filename component."""
    return name.replace(" ", "_").replace("/", "-").replace("\\", "-")


def _img_ext(region: dict) -> str:
    """Extract image extension from a region dict, defaulting to .jpg."""
    orig = region.get("image", {}).get("data", {}).get("filename", "")
    ext = Path(orig).suffix if orig else ".jpg"
    return ext or ".jpg"


def _save_image(client, url: str, filepath: Path) -> bool:
    """Download and save an image. Returns True on success."""
    if filepath.exists():
        logger.info("  Already exists: %s", filepath)
        return True
    try:
        filepath.write_bytes(client.download_image(url))
        logger.info("  Downloaded: %s", filepath)
        return True
    except Exception as exc:
        logger.warning("  Failed to download %s: %s", filepath.name, exc)
        return False


def _save_cropped_image(client, url: str, bbox: dict, filepath: Path) -> bool:
    """Download a full image, crop it to *bbox*, and save the crop."""
    if filepath.exists():
        logger.info("  Already exists: %s", filepath)
        return True
    try:
        img_data = client.download_image(url)
        img = PILImage.open(io.BytesIO(img_data))
        w, h = img.size
        left = int(bbox["xmin"] * w)
        upper = int(bbox["ymin"] * h)
        right = int(bbox["xmax"] * w)
        lower = int(bbox["ymax"] * h)
        crop = img.crop((left, upper, right, lower))
        filepath = filepath.with_suffix(".png")
        crop.save(filepath)
        logger.info("  Cropped & saved: %s (%dx%d)", filepath, right - left, lower - upper)
        return True
    except Exception as exc:
        logger.warning("  Failed to crop %s: %s", filepath.name, exc)
        return False


def _draw_bboxes(img: PILImage.Image, bboxes: list[dict], color=(255, 107, 53), thickness: int = 3) -> PILImage.Image:
    """Draw normalized bboxes on a PIL image. Returns a copy with overlays."""
    from PIL import ImageDraw
    img = img.copy()
    draw = ImageDraw.Draw(img)
    w, h = img.size
    for bbox in bboxes:
        x0 = int(bbox["xmin"] * w)
        y0 = int(bbox["ymin"] * h)
        x1 = int(bbox["xmax"] * w)
        y1 = int(bbox["ymax"] * h)
        for i in range(thickness):
            draw.rectangle([x0 - i, y0 - i, x1 + i, y1 + i], outline=color)
    return img


def _download_n_per_concept(client, view_id, view_label, tag_ids, concept_map, view_dir, *, n=2, crop=False):
    """TAG/CLA: fetch N images per concept tag."""
    downloaded = 0
    for tag_id in tag_ids:
        concept_name = _dl_sanitize(concept_map.get(tag_id, str(tag_id)))
        try:
            regions = client.get_regions(view_id, page_size=n, tag=tag_id)
        except Exception as exc:
            logger.warning("  Could not fetch regions for %s / %s: %s", view_label, concept_name, exc)
            continue
        if not regions:
            logger.info("  No image found for %s / %s", view_label, concept_name)
            continue
        for idx, region in enumerate(regions[:n], 1):
            img_url = region.get("image", {}).get("original_signed_url")
            if not img_url:
                continue
            bbox = region.get("region", {}).get("bbox") if crop else None
            ext = _img_ext(region)
            filepath = view_dir / f"{view_label}__{concept_name}__{idx}{ext}"
            if bbox:
                ok = _save_cropped_image(client, img_url, bbox, filepath)
            else:
                ok = _save_image(client, img_url, filepath)
            if ok:
                downloaded += 1
    if downloaded == 0:
        _download_fallback(client, view_id, view_label, view_dir, count=1)


def _download_det_per_concept(client, view_id, view_label, tag_ids, concept_map, view_dir, *, n=2):
    """DET: fetch N images per concept with bbox overlays."""
    DET_BBOX_COLOR = (255, 107, 53)
    downloaded = 0
    for tag_id in tag_ids:
        concept_name = _dl_sanitize(concept_map.get(tag_id, str(tag_id)))
        try:
            regions = client.get_regions(view_id, page_size=n, tag=tag_id)
        except Exception as exc:
            logger.warning("  Could not fetch regions for DET %s / %s: %s", view_label, concept_name, exc)
            continue
        if not regions:
            continue
        for idx, region in enumerate(regions[:n], 1):
            img_url = region.get("image", {}).get("original_signed_url")
            region_id = region.get("region", {}).get("id")
            if not img_url or not region_id:
                continue
            filepath = view_dir / f"{view_label}__{concept_name}__{idx}.png"
            if filepath.exists():
                downloaded += 1
                continue
            try:
                annotations = client.get_annotations(view_id, region_id)
            except Exception:
                _save_image(client, img_url, filepath)
                downloaded += 1
                continue
            bboxes = []
            for ann in annotations:
                ann_tags = ann.get("tags", [])
                tag_match = any(
                    (t if isinstance(t, int) else t.get("id")) == tag_id
                    for t in ann_tags
                )
                if tag_match:
                    bbox = ann.get("region", {}).get("bbox")
                    if bbox:
                        bboxes.append(bbox)
            try:
                img_data = client.download_image(img_url)
                img = PILImage.open(io.BytesIO(img_data))
                if bboxes:
                    img = _draw_bboxes(img, bboxes, color=DET_BBOX_COLOR, thickness=4)
                img.save(filepath)
                downloaded += 1
            except Exception as exc:
                logger.warning("  Failed to save DET image %s: %s", filepath.name, exc)
    if downloaded == 0:
        _download_fallback(client, view_id, view_label, view_dir, count=1)


def _download_fallback(client, view_id, view_label, view_dir, count=1):
    """Download generic sample images when no tag-specific strategy works."""
    try:
        regions = client.get_regions(view_id, page_size=count)
    except Exception:
        return
    for i, region in enumerate(regions[:count]):
        img_url = region.get("image", {}).get("original_signed_url")
        if not img_url:
            continue
        ext = _img_ext(region)
        filepath = view_dir / f"{view_label}__sample_{i + 1}{ext}"
        _save_image(client, img_url, filepath)


def download_sample_images(client, project_map: dict, images_dir: Path):
    """Download sample images for each view, organised by view name.

    Strategy per view kind:
    - TAG / CLA: 2 images per concept.
    - DET: 2 images per concept with bbox overlays.
    """
    images_dir.mkdir(parents=True, exist_ok=True)
    concept_map = {c["id"]: c["concept_name"] for c in project_map.get("concepts", [])}

    node_map = {n["id"]: n for n in project_map["nodes"]}
    parent_kind: dict[str, str] = {}
    for edge in project_map.get("edges", []):
        src = node_map.get(edge["source"], {})
        parent_kind[edge["target"]] = src.get("data", {}).get("kind", "").upper()

    for node in project_map["nodes"]:
        view_id = node["id"]
        view_label = _dl_sanitize(node["label"])
        kind = node["data"].get("kind", "").upper()
        tag_ids: list[int] = node["data"].get("tag_ids", [])
        is_child_of_det = parent_kind.get(view_id) == "DET"

        view_dir = images_dir / view_label
        view_dir.mkdir(parents=True, exist_ok=True)

        if kind in ("TAG", "CLA"):
            _download_n_per_concept(client, view_id, view_label, tag_ids,
                                    concept_map, view_dir, n=2, crop=is_child_of_det)
        elif kind == "DET":
            _download_det_per_concept(client, view_id, view_label, tag_ids,
                                      concept_map, view_dir, n=2)
        else:
            _download_fallback(client, view_id, view_label, view_dir, count=2)

    logger.info("Images saved to %s", images_dir)
