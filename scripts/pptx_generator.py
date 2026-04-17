"""
Generate a skeleton .pptx annotation guide from a Deepomatic project map.

Adapted from Brieuc's PR #317 in deepomatic-tools.

The project map JSON has the following structure:
  {
    "nodes": [ { "id": ..., "label": ..., "data": { "parent": ..., "kind": ..., "conditions": ... } } ],
    "edges": [ { "source": ..., "target": ..., "data": {} } ],
    "concepts": [ { "id": ..., "concept_name": ... } ]
  }

View kinds:
  CLA  – Classification
  DET  – Detection  (bounding-box based)
  TAG  – Tagging
"""

import logging
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Pt

logger = logging.getLogger(__name__)


def _sanitize(name: str) -> str:
    """Turn a label into a safe directory/filename component."""
    return name.replace(" ", "_").replace("/", "-").replace("\\", "-")

# ---------------------------------------------------------------------------
# Slide-layout indices (default Office Theme)
# ---------------------------------------------------------------------------
TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "template" / "template.pptx"

# Layout indices in template/template.pptx (see the template's slide master)
LAYOUT_COVER = 0          # "Title Only Slide"       — big centred title for cover page
LAYOUT_SECTION = 5        # "Title + Subtitle slide" — section divider between views
LAYOUT_CONTENT = 11       # "Title and Content"      — info slides with bullet content
LAYOUT_COMPARISON = 6     # "Comparison"             — two-column content
LAYOUT_IMAGES = 13        # "Title Only"             — concept image slides

# Kind labels
KIND_LABELS = {
    "CLA": "Classification",
    "DET": "Detection",
    "TAG": "Tagging",
}

# Colours
COLOUR_TITLE = RGBColor(0x1F, 0x49, 0x7D)  # dark blue
COLOUR_DET_BOX = RGBColor(0xFF, 0x6F, 0x00)  # orange
COLOUR_CLA_BOX = RGBColor(0x70, 0xAD, 0x47)  # green
COLOUR_TAG_BOX = RGBColor(0x5B, 0x9B, 0xD5)  # blue
COLOUR_PLACEHOLDER_BG = RGBColor(0xF2, 0xF2, 0xF2)  # light grey


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _clear_template_slides(prs):
    """Drop any sample slides shipped with the template — we only want its masters/layouts."""
    sld_id_lst = prs.slides._sldIdLst
    rid_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(rid_attr))
        sld_id_lst.remove(sld_id)


def _set_text(tf, text, bold=False, size_pt=14, colour=None, align=PP_ALIGN.LEFT):
    """Replace all text in a text-frame with a single run."""
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size_pt)
    if colour:
        run.font.color.rgb = colour


def _add_text_box(slide, left, top, width, height, text, bold=False, size_pt=12, colour=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_text(tf, text, bold=bold, size_pt=size_pt, colour=colour, align=align)
    return txBox


def _add_placeholder_image(slide, left, top, width, height, label, box_colour):
    """Draw a labelled rectangle that acts as a picture placeholder."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOUR_PLACEHOLDER_BG
    shape.line.color.rgb = box_colour
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = box_colour
    run.font.bold = True
    return shape


def _add_bbox_overlay(slide, parent_left, parent_top, parent_width, parent_height, box_colour):
    """Draw a small bounding-box overlay inside the image placeholder."""
    margin_w = parent_width * 0.15
    margin_h = parent_height * 0.15
    box = slide.shapes.add_shape(
        1,
        parent_left + margin_w,
        parent_top + margin_h,
        parent_width - 2 * margin_w,
        parent_height - 2 * margin_h,
    )
    box.fill.background()
    box.line.color.rgb = box_colour
    box.line.width = Pt(2)
    return box


# ---------------------------------------------------------------------------
# Map parsing
# ---------------------------------------------------------------------------


def _build_tree(project_map):
    """
    Return (nodes_dict, roots_list) from a project map.
    Each node: { "id", "label", "kind", "parent", "conditions", "children" }
    """
    nodes = {}
    for n in project_map["nodes"]:
        nodes[n["id"]] = {
            "id": n["id"],
            "label": n["label"],
            "kind": n["data"].get("kind", ""),
            "parent": n["data"].get("parent", ""),
            "conditions": n["data"].get("conditions") or [],
            "tag_names": n["data"].get("tag_names", []),
            "children": [],
        }

    for edge in project_map["edges"]:
        src, tgt = edge["source"], edge["target"]
        if src and src in nodes and tgt in nodes:
            nodes[src]["children"].append(tgt)

    roots = [nid for nid, n in nodes.items() if not n["parent"]]
    return nodes, roots


def _build_concept_map(project_map):
    """Return dict {concept_id: concept_name}."""
    return {c["id"]: c["concept_name"] for c in project_map.get("concepts", [])}


def _resolve_conditions(conditions, concept_map):
    """Turn condition lists into a human-readable string."""
    if not conditions:
        return ""
    groups = []
    for group in conditions:
        names = [concept_map.get(cid, str(cid)) for cid in group]
        groups.append(" & ".join(names))
    return "  |  ".join(groups)


def _dfs_order(nodes, roots):
    """Return node IDs in DFS order starting from roots."""
    order = []
    visited = set()

    def dfs(nid):
        if nid in visited:
            return
        visited.add(nid)
        order.append(nid)
        for child_id in nodes[nid]["children"]:
            dfs(child_id)

    for r in roots:
        dfs(r)
    for nid in nodes:
        if nid not in visited:
            order.append(nid)
    return order


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _kind_colour(kind):
    if kind == "DET":
        return COLOUR_DET_BOX
    if kind == "CLA":
        return COLOUR_CLA_BOX
    return COLOUR_TAG_BOX


def _find_view_images(view_label: str, images_dir) -> dict[str, Path | None]:
    """Return a dict {concept_name: image_path_or_None} for a view.

    Images are expected at  images_dir/<sanitized_view>/<view>__<concept>.ext
    """
    if images_dir is None:
        return {}
    view_dir = Path(images_dir) / _sanitize(view_label)
    if not view_dir.is_dir():
        return {}
    exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}
    result: dict[str, Path] = {}
    for p in sorted(view_dir.iterdir()):
        if p.suffix.lower() not in exts:
            continue
        stem = p.stem
        if "__" in stem:
            concept = stem.split("__", 1)[1].replace("_", " ")
        else:
            concept = stem
        result[concept] = p
    return result


def _add_concept_to_slide(
    slide, concept_name: str, img_path: Path | None,
    x, y, col_w, img_area_h, box_colour, slide_w, slide_h,
):
    """Place one concept (image + legend) in a column on the slide."""
    from PIL import Image as PILImage

    legend_h = Cm(1.2)

    if img_path and img_path.exists():
        try:
            with PILImage.open(img_path) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = 0, 0

        if iw > 0 and ih > 0:
            # Size constraints: max 1/3, min 1/4 of max(slide_w, slide_h)
            max_dim = max(slide_w, slide_h)
            size_max = max_dim // 3
            size_min = max_dim // 4

            # Fit within column and image area
            max_w = min(col_w, size_max)
            max_h = min(img_area_h - legend_h, size_max)

            scale = min(max_w / iw, max_h / ih)
            pic_w = int(iw * scale)
            pic_h = int(ih * scale)

            # Enforce minimum size
            if max(pic_w, pic_h) < size_min:
                upscale = size_min / max(pic_w, pic_h)
                pic_w = int(pic_w * upscale)
                pic_h = int(pic_h * upscale)

            # Centre horizontally in column
            x_off = (col_w - pic_w) // 2
            # Place image vertically centred in the available area
            y_off = (img_area_h - legend_h - pic_h) // 2

            try:
                slide.shapes.add_picture(
                    str(img_path), x + x_off, y + max(y_off, 0), pic_w, pic_h,
                )
            except Exception as exc:
                logger.warning("Could not embed %s: %s", img_path.name, exc)
    else:
        # No image — draw a light placeholder rectangle
        margin = Cm(1)
        placeholder_w = col_w - 2 * margin
        placeholder_h = img_area_h - legend_h - 2 * margin
        if placeholder_w > 0 and placeholder_h > 0:
            shape = slide.shapes.add_shape(
                1, x + margin, y + margin, placeholder_w, placeholder_h,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOUR_PLACEHOLDER_BG
            shape.line.color.rgb = box_colour
            shape.line.width = Pt(1)
            shape.line.dash_style = 2  # dash
            tf = shape.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = "[ image missing ]"
            run.font.size = Pt(10)
            run.font.color.rgb = box_colour

    # Legend — concept name centred below the image area
    _add_text_box(
        slide, x, y + img_area_h - legend_h, col_w, legend_h,
        concept_name, bold=True, size_pt=12, colour=box_colour, align=PP_ALIGN.CENTER,
    )


def _slides_for_view(prs, node, nodes, concept_map, images_dir=None):
    """Create one info slide + concept slides (max 3 concepts each) for a view."""
    kind = node["kind"]
    kind_label = KIND_LABELS.get(kind, kind)
    box_colour = _kind_colour(kind)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # --- Info slide (always): layout 11 "Title and Content" — fill the content placeholder ---
    info_layout = prs.slide_layouts[LAYOUT_CONTENT]
    slide = prs.slides.add_slide(info_layout)

    slide.shapes.title.text = f"{node['label']}  –  {kind_label}"

    parent_label = ""
    if node["parent"] and node["parent"] in nodes:
        parent_label = nodes[node["parent"]]["label"]
    cond_str = _resolve_conditions(node["conditions"], concept_map)
    child_labels = [nodes[c]["label"] for c in node["children"] if c in nodes]
    tag_names = node.get("tag_names", [])

    lines = [
        ("Parent view", parent_label or "— (root)"),
        ("Activated by", cond_str or "—"),
        ("Child views", ", ".join(child_labels) if child_labels else "—"),
    ]
    if tag_names:
        lines.append(("Concepts", ", ".join(tag_names)))

    content_ph = slide.placeholders[1]  # "Content Placeholder 2"
    tf = content_ph.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (label, value) in enumerate(lines):
        # Label line at level 0 (inherits template bullet style)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = 0
        label_run = para.add_run()
        label_run.text = label
        label_run.font.bold = True
        # Value line at level 1 (indented, template sub-bullet style)
        val_para = tf.add_paragraph()
        val_para.level = 1
        val_run = val_para.add_run()
        val_run.text = value

    # --- Concept slides (max 3 per slide): layout 13 "Title Only" ---
    concept_layout = prs.slide_layouts[LAYOUT_IMAGES]
    view_images = _find_view_images(node["label"], images_dir)

    # Build the list of concepts to show
    concepts: list[str] = list(tag_names) if tag_names else []

    # If no concepts defined but we have images, use image names
    if not concepts and view_images:
        concepts = list(view_images.keys())

    # If still nothing, no concept slides needed (info slide is enough)
    if not concepts:
        return

    # Chunk concepts into groups of 3
    CONCEPTS_PER_SLIDE = 3
    chunks = [concepts[i:i + CONCEPTS_PER_SLIDE] for i in range(0, len(concepts), CONCEPTS_PER_SLIDE)]

    # Content zone matches layout 11's content placeholder so info + concept slides align.
    content_left = Cm(3.6)
    content_top = Cm(5.0)
    content_w = Cm(28.0)
    content_h = slide_h - content_top - Cm(1.5)

    for chunk_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(concept_layout)

        title = slide.shapes.title
        if len(chunks) > 1:
            title.text = f"{node['label']}  ({chunk_idx + 1}/{len(chunks)})"
        else:
            title.text = node["label"]

        n_cols = len(chunk)
        col_w = content_w // n_cols

        for col_idx, concept_name in enumerate(chunk):
            # Try to find the image — match by concept name
            img_path = view_images.get(concept_name)
            # Also try sanitized matching
            if img_path is None:
                sanitized = _sanitize(concept_name).replace("_", " ")
                for k, v in view_images.items():
                    if k == sanitized or k.lower() == concept_name.lower():
                        img_path = v
                        break
            # Also try if image key contains the concept (partial match for "sample" names)
            if img_path is None:
                cn_lower = concept_name.lower().replace(" ", "_")
                for k, v in view_images.items():
                    if cn_lower in k.lower().replace(" ", "_"):
                        img_path = v
                        break

            x = content_left + col_idx * col_w
            _add_concept_to_slide(
                slide, concept_name, img_path,
                x, content_top, col_w, content_h,
                box_colour, slide_w, slide_h,
            )


def _slide_intro(prs, nodes, roots, concept_map, *, org_slug="", project_slug=""):
    """Cover slide + tree diagram overview slide."""
    # --- Cover slide (Layout 0 — big centred title) ---
    cover = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])
    title_parts = ["Annotation Guide"]
    if org_slug or project_slug:
        title_parts.append(" · ")
        title_parts.append(" - ".join(filter(None, [org_slug.upper(), project_slug.upper()])))
    cover.shapes.title.text = "".join(title_parts)

    # --- Tree diagram slide ---
    _slide_tree_overview(prs, nodes, roots, concept_map)


# ---------------------------------------------------------------------------
# Tree diagram
# ---------------------------------------------------------------------------

def _compute_tree_positions(nodes, roots):
    """Return dict {nid: (grid_x, grid_y)} and (total_width, num_levels)."""
    widths: dict[str, float] = {}

    def calc_width(nid):
        children = [c for c in nodes[nid]["children"] if c in nodes]
        if not children:
            widths[nid] = 1.0
        else:
            widths[nid] = sum(calc_width(c) for c in children)
        return widths[nid]

    total_w = sum(calc_width(r) for r in roots)

    positions: dict[str, tuple[float, int]] = {}
    max_depth = 0

    def assign(nid, x_start, depth):
        nonlocal max_depth
        max_depth = max(max_depth, depth)
        w = widths[nid]
        positions[nid] = (x_start + w / 2.0, depth)
        children = [c for c in nodes[nid]["children"] if c in nodes]
        cx = x_start
        for c in children:
            assign(c, cx, depth + 1)
            cx += widths[c]

    x = 0.0
    for r in roots:
        assign(r, x, 0)
        x += widths[r]

    return positions, total_w, max_depth + 1


def _slide_tree_overview(prs, nodes, roots, concept_map):
    """Draw a tree diagram of the view hierarchy on a slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_IMAGES])
    slide.shapes.title.text = "Views Overview"

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    positions, total_w, num_levels = _compute_tree_positions(nodes, roots)
    if not positions:
        return slide

    # Drawing area (below title)
    area_left = Cm(1.5)
    area_top = Cm(2.2)
    area_w = slide_w - Cm(3.0)
    area_h = slide_h - area_top - Cm(0.8)

    # Node box dimensions – scale to fit
    node_w = min(Cm(3.2), int(area_w / max(total_w, 1) - Cm(0.4)))
    node_w = max(node_w, Cm(1.8))
    node_h = Cm(1.0)

    # Grid cell sizes
    cell_w = area_w / total_w if total_w else area_w
    cell_h = area_h / num_levels if num_levels > 1 else area_h

    # Convert grid → EMU positions (centre of each node box)
    emu_pos: dict[str, tuple[int, int]] = {}
    for nid, (gx, gy) in positions.items():
        cx = int(area_left + gx * cell_w)
        cy = int(area_top + gy * cell_h + node_h // 2)
        emu_pos[nid] = (cx, cy)

    # --- Draw connectors (edges) ---
    connector_colour = RGBColor(0xAA, 0xAA, 0xAA)
    cond_label_colour = RGBColor(0x55, 0x55, 0x55)

    for nid in positions:
        children = [c for c in nodes[nid]["children"] if c in nodes]
        if not children:
            continue
        px, py = emu_pos[nid]
        parent_bottom_y = py + node_h // 2

        for child_id in children:
            cx, cy = emu_pos[child_id]
            child_top_y = cy - node_h // 2

            # Straight line from parent bottom-centre → child top-centre
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                px, parent_bottom_y, cx, child_top_y,
            )
            conn.line.color.rgb = connector_colour
            conn.line.width = Pt(1.0)

            # Condition label at midpoint of the edge
            cond_text = _resolve_conditions(nodes[child_id]["conditions"], concept_map)
            if cond_text:
                mid_x = (px + cx) // 2
                mid_y = (parent_bottom_y + child_top_y) // 2
                lbl_w = min(Cm(4.0), int(cell_w * 0.95))
                lbl_h = Cm(0.9)
                tb = slide.shapes.add_textbox(
                    mid_x - lbl_w // 2, mid_y - lbl_h // 2, lbl_w, lbl_h,
                )
                tf = tb.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                run = para.add_run()
                run.text = f"({cond_text})"
                run.font.size = Pt(6)
                run.font.color.rgb = cond_label_colour

    # --- Draw node boxes ---
    for nid, (cx, cy) in emu_pos.items():
        n = nodes[nid]
        colour = _kind_colour(n["kind"])

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            cx - node_w // 2, cy - node_h // 2, node_w, node_h,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.color.rgb = colour
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = n["label"]
        run.font.size = Pt(8)
        run.font.bold = True

    return slide


def _slide_section(prs, node):
    """Add a section divider slide (Layout 5) for a top-level view."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    slide.shapes.title.text = node["label"]
    kind_label = KIND_LABELS.get(node["kind"], node["kind"])
    # Layout 5: placeholder 0 = title (bottom), placeholder 1 = subtitle (above)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = kind_label
            break
    return slide


# _slide_for_view replaced by _slides_for_view above


# ---------------------------------------------------------------------------
# Main generator
# ---------------------------------------------------------------------------


def generate_pptx(project_map: dict, output_path: str | Path, *, images_dir: Path | None = None, org_slug: str = "", project_slug: str = "") -> Path:
    """
    Build the annotation guide .pptx from a project map and write it.

    Parameters
    ----------
    project_map : dict
        Parsed JSON project map containing "nodes", "edges", "concepts".
    output_path : str | Path
        Destination file path.
    images_dir : Path | None
        Directory containing downloaded sample images (one sub-folder per view).
    org_slug : str
        Organisation slug (shown on the cover slide).
    project_slug : str
        Project slug (shown on the cover slide).
    """
    nodes, roots = _build_tree(project_map)
    concept_map = _build_concept_map(project_map)
    ordered = _dfs_order(nodes, roots)

    prs = Presentation(str(TEMPLATE_PATH))
    _clear_template_slides(prs)
    logger.info("Using template at %s", TEMPLATE_PATH)

    _slide_intro(prs, nodes, roots, concept_map, org_slug=org_slug, project_slug=project_slug)

    for nid in ordered:
        node = nodes[nid]
        # Section divider for root-level views (no parent)
        if not node["parent"] or node["parent"] not in nodes:
            _slide_section(prs, node)
        _slides_for_view(prs, node, nodes, concept_map, images_dir=images_dir)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("Saved annotation guide to %s  (%d slides)", output_path, len(prs.slides))
    return output_path
