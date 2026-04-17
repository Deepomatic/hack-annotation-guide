"""
Slide composition for the Annotation Guide.

This file defines what slides appear in the deck and in what order.
Edit this file to add, remove, or reorder slides.
Add hardcoded content slides by writing new build_xxx() functions
and calling them from build_all_slides().

All PPTX primitives come from pptx_helper — do NOT use raw python-pptx here.
"""

import logging
import math
from collections import defaultdict
from pathlib import Path

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Emu, Pt

from pptx_helper import (
    # Colors
    NAVY,
    NAVY_LIGHT,
    WHITE,
    LIGHT_BG,
    MUTED,
    DIVIDER,
    DARK_TEXT,
    ORANGE,
    TEAL,
    SKY_BLUE,
    GREEN,
    RED,
    PLACEHOLDER_BG,
    # Typography
    FONT_FAMILY,
    FONT_SIZE_COVER_TITLE,
    FONT_SIZE_COVER_SUBTITLE,
    FONT_SIZE_SECTION_TITLE,
    FONT_SIZE_SECTION_SUBTITLE,
    FONT_SIZE_SLIDE_TITLE,
    FONT_SIZE_BODY,
    FONT_SIZE_BODY_SMALL,
    FONT_SIZE_CAPTION,
    FONT_SIZE_LABEL,
    FONT_SIZE_BADGE,
    # Layout
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
    MARGIN_LEFT,
    MARGIN_RIGHT,
    MARGIN_TOP,
    MARGIN_BOTTOM,
    CONTENT_LEFT,
    CONTENT_TOP,
    CONTENT_WIDTH,
    CONTENT_HEIGHT,
    TITLE_LEFT,
    TITLE_TOP,
    TITLE_WIDTH,
    TITLE_HEIGHT,
    ACCENT_BAR_WIDTH,
    ACCENT_BAR_HEIGHT,
    # Functions
    add_blank_slide,
    set_slide_bg_solid,
    set_slide_bg_gradient,
    add_textbox,
    add_rich_textbox,
    add_multiline_textbox,
    add_title,
    add_subtitle,
    add_rounded_rect,
    add_accent_bar,
    add_line,
    add_badge,
    add_image,
    add_image_placeholder,
    grid_positions,
    add_tree_connector,
    add_footer,
    add_bottom_strip,
    add_top_line,
)

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────────────────────────────
# Kind constants
# ──────────────────────────────────────────────────────────────────────

KIND_LABELS = {
    "CLA": "Classification",
    "DET": "Detection",
    "TAG": "Tagging",
}


def kind_color(kind: str):
    """Return the accent color for a view kind."""
    if kind == "DET":
        return ORANGE
    if kind == "CLA":
        return TEAL
    return SKY_BLUE


# ──────────────────────────────────────────────────────────────────────
# Project map parsing helpers
# ──────────────────────────────────────────────────────────────────────


def _sanitize(name: str) -> str:
    return name.replace(" ", "_").replace("/", "-").replace("\\", "-")


def build_tree(project_map: dict):
    """Parse project map into (nodes_dict, roots_list).

    Each node: {id, label, kind, parent, conditions, tag_names, children}
    """
    nodes = {}
    for n in project_map["nodes"]:
        nodes[n["id"]] = {
            "id": n["id"],
            "label": n["label"],
            "kind": n["data"].get("kind", ""),
            "parent": n["data"].get("parent", ""),
            "conditions": n["data"].get("conditions") or [],
            "tag_names": n["data"].get("tag_names", []),
            "children": [],
        }

    for edge in project_map["edges"]:
        src, tgt = edge["source"], edge["target"]
        if src and src in nodes and tgt in nodes:
            nodes[src]["children"].append(tgt)

    roots = [nid for nid, n in nodes.items() if not n["parent"]]
    return nodes, roots


def build_concept_map(project_map: dict) -> dict:
    return {c["id"]: c["concept_name"] for c in project_map.get("concepts", [])}


def resolve_conditions(conditions, concept_map: dict) -> str:
    if not conditions:
        return ""
    groups = []
    for group in conditions:
        names = [concept_map.get(cid, str(cid)) for cid in group]
        groups.append(" & ".join(names))
    return "  |  ".join(groups)


def dfs_order(nodes: dict, roots: list) -> list[str]:
    """Return node IDs in depth-first order."""
    order = []
    visited = set()

    def dfs(nid):
        if nid in visited:
            return
        visited.add(nid)
        order.append(nid)
        for child_id in nodes[nid]["children"]:
            dfs(child_id)

    for r in roots:
        dfs(r)
    # Catch orphans
    for nid in nodes:
        if nid not in visited:
            order.append(nid)
    return order


def find_view_images(view_label: str, images_dir) -> dict[str, Path | None]:
    """Return {concept_name: image_path} for a view's downloaded images."""
    if images_dir is None:
        return {}
    view_dir = Path(images_dir) / _sanitize(view_label)
    if not view_dir.is_dir():
        return {}
    exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}
    result: dict[str, Path] = {}
    for p in sorted(view_dir.iterdir()):
        if p.suffix.lower() not in exts:
            continue
        stem = p.stem
        if "__" in stem:
            concept = stem.split("__", 1)[1].replace("_", " ")
        else:
            concept = stem
        result[concept] = p
    return result


# ──────────────────────────────────────────────────────────────────────
# Tree layout computation
# ──────────────────────────────────────────────────────────────────────


def _compute_tree_positions(nodes, roots):
    """Return {nid: (grid_x, grid_y)}, total_width, num_levels."""
    widths: dict[str, float] = {}

    def calc_width(nid):
        children = [c for c in nodes[nid]["children"] if c in nodes]
        if not children:
            widths[nid] = 1.0
        else:
            widths[nid] = sum(calc_width(c) for c in children)
        return widths[nid]

    total_w = sum(calc_width(r) for r in roots)

    positions: dict[str, tuple[float, int]] = {}
    max_depth = 0

    def assign(nid, x_start, depth):
        nonlocal max_depth
        max_depth = max(max_depth, depth)
        w = widths[nid]
        positions[nid] = (x_start + w / 2.0, depth)
        children = [c for c in nodes[nid]["children"] if c in nodes]
        cx = x_start
        for c in children:
            assign(c, cx, depth + 1)
            cx += widths[c]

    x = 0.0
    for r in roots:
        assign(r, x, 0)
        x += widths[r]

    return positions, total_w, max_depth + 1


# ══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS — edit / add / reorder these as needed
# ══════════════════════════════════════════════════════════════════════


def build_cover_slide(prs, project_name: str = ""):
    """Slide 1: Full navy cover with project name."""
    slide = add_blank_slide(prs)
    set_slide_bg_gradient(slide, NAVY, NAVY_LIGHT)

    # Decorative accent bar at top
    add_accent_bar(slide, 0, 0, SLIDE_WIDTH, Cm(0.25), ORANGE)

    # Title — centered vertically
    add_textbox(
        slide,
        MARGIN_LEFT, Cm(6.0),
        SLIDE_WIDTH - 2 * MARGIN_LEFT, Cm(3.0),
        "Annotation Guide",
        font_size=FONT_SIZE_COVER_TITLE,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Project name subtitle
    if project_name:
        add_textbox(
            slide,
            MARGIN_LEFT, Cm(9.5),
            SLIDE_WIDTH - 2 * MARGIN_LEFT, Cm(2.0),
            project_name,
            font_size=FONT_SIZE_COVER_SUBTITLE,
            color=MUTED,
            alignment=PP_ALIGN.CENTER,
        )

    # Bottom strip
    add_accent_bar(slide, 0, SLIDE_HEIGHT - Cm(0.25), SLIDE_WIDTH, Cm(0.25), ORANGE)

    return slide


def build_toc_slide(prs, nodes: dict, roots: list):
    """Slide 2 (optional): Table of Contents listing root views."""
    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, WHITE)

    add_title(slide, "Table of Contents")
    add_top_line(slide)

    y = CONTENT_TOP + Cm(0.5)
    for i, root_id in enumerate(roots):
        node = nodes[root_id]
        accent = kind_color(node["kind"])
        kind_lbl = KIND_LABELS.get(node["kind"], node["kind"])

        # Accent dot
        add_accent_bar(slide, MARGIN_LEFT, y + Cm(0.25), Cm(0.3), Cm(0.3), accent)

        # View name
        add_textbox(
            slide,
            MARGIN_LEFT + Cm(0.8), y,
            Cm(20), Cm(0.8),
            f"{node['label']}",
            font_size=FONT_SIZE_BODY,
            bold=True,
            color=DARK_TEXT,
        )

        # Kind label
        add_textbox(
            slide,
            MARGIN_LEFT + Cm(21), y,
            Cm(8), Cm(0.8),
            kind_lbl,
            font_size=FONT_SIZE_BODY_SMALL,
            color=MUTED,
        )

        y += Cm(1.2)

    add_bottom_strip(slide, NAVY)
    return slide


def build_overview_slide(prs, nodes: dict, roots: list, concept_map: dict):
    """Tree diagram overview of the view hierarchy."""
    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, WHITE)

    add_title(slide, "Views Overview")
    add_top_line(slide)

    positions, total_w, num_levels = _compute_tree_positions(nodes, roots)
    if not positions:
        return slide

    area_left = Cm(1.5)
    area_top = Cm(4.5)
    area_w = SLIDE_WIDTH - Cm(3.0)
    area_h = SLIDE_HEIGHT - area_top - Cm(1.5)

    # Node box dimensions
    node_w = min(Cm(4.0), int(int(area_w) / max(total_w, 1) - int(Cm(0.6))))
    node_w = max(node_w, Cm(2.0))
    node_h = Cm(1.2)

    cell_w = int(area_w) / total_w if total_w else int(area_w)
    cell_h = int(area_h) / num_levels if num_levels > 1 else int(area_h)

    # Convert grid → pixel positions
    emu_pos: dict[str, tuple[int, int]] = {}
    for nid, (gx, gy) in positions.items():
        cx = int(int(area_left) + gx * cell_w)
        cy = int(int(area_top) + gy * cell_h + int(node_h) // 2)
        emu_pos[nid] = (cx, cy)

    # Draw connectors first (below nodes)
    for nid in positions:
        children = [c for c in nodes[nid]["children"] if c in nodes]
        if not children:
            continue
        px, py = emu_pos[nid]
        parent_bottom_y = py + int(node_h) // 2

        for child_id in children:
            ccx, ccy = emu_pos[child_id]
            child_top_y = ccy - int(node_h) // 2

            add_tree_connector(
                slide, px, parent_bottom_y, ccx, child_top_y,
                color=DIVIDER, width=Pt(1.5),
            )

            # Condition label
            cond_text = resolve_conditions(nodes[child_id]["conditions"], concept_map)
            if cond_text:
                mid_x = (px + ccx) // 2
                mid_y = (parent_bottom_y + child_top_y) // 2
                lbl_w = min(Cm(5.0), int(cell_w * 0.95))
                add_textbox(
                    slide,
                    mid_x - lbl_w // 2, mid_y - Cm(0.4),
                    lbl_w, Cm(0.8),
                    f"({cond_text})",
                    font_size=Pt(7),
                    color=MUTED,
                    alignment=PP_ALIGN.CENTER,
                )

    # Draw node boxes
    for nid, (cx, cy) in emu_pos.items():
        n = nodes[nid]
        accent = kind_color(n["kind"])

        add_rounded_rect(
            slide,
            cx - int(node_w) // 2, cy - int(node_h) // 2,
            node_w, node_h,
            fill_color=WHITE,
            border_color=accent,
            border_width=Pt(2),
            text=n["label"],
            text_color=DARK_TEXT,
            text_size=Pt(9),
            text_bold=True,
        )

    # Legend at bottom
    legend_y = SLIDE_HEIGHT - Cm(1.8)
    legend_x = MARGIN_LEFT
    for kind_code, kind_label in KIND_LABELS.items():
        c = kind_color(kind_code)
        add_accent_bar(slide, legend_x, legend_y + Cm(0.15), Cm(0.5), Cm(0.5), c)
        add_textbox(
            slide,
            legend_x + Cm(0.7), legend_y,
            Cm(4), Cm(0.8),
            kind_label,
            font_size=FONT_SIZE_LABEL,
            color=MUTED,
        )
        legend_x += Cm(5)

    add_bottom_strip(slide, NAVY)
    return slide


def build_section_slide(prs, view_name: str, view_kind: str, *, accent_color=None):
    """Section divider slide for a root-level view. Navy gradient background."""
    slide = add_blank_slide(prs)
    set_slide_bg_gradient(slide, NAVY, NAVY_LIGHT)

    accent = accent_color or kind_color(view_kind)

    # Accent bar left side
    add_accent_bar(slide, 0, Cm(5), ACCENT_BAR_WIDTH, Cm(9), accent)

    # View name — large
    add_textbox(
        slide,
        Cm(3.0), Cm(6.0),
        Cm(26), Cm(3.0),
        view_name,
        font_size=FONT_SIZE_SECTION_TITLE,
        bold=True,
        color=WHITE,
    )

    # Kind badge
    kind_label = KIND_LABELS.get(view_kind, view_kind)
    add_badge(
        slide,
        Cm(3.0), Cm(9.5),
        kind_label,
        bg_color=accent,
        text_color=WHITE,
    )

    # Bottom strip
    add_accent_bar(slide, 0, SLIDE_HEIGHT - Cm(0.25), SLIDE_WIDTH, Cm(0.25), accent)

    return slide


def build_info_slide(prs, node: dict, nodes: dict, concept_map: dict):
    """Info slide for a view: key-value metadata, accent bar, clean layout."""
    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, WHITE)

    kind = node["kind"]
    accent = kind_color(kind)

    # Title
    add_title(slide, node["label"])

    # Kind badge next to title
    kind_label = KIND_LABELS.get(kind, kind)
    add_badge(
        slide,
        SLIDE_WIDTH - MARGIN_RIGHT - Cm(5), TITLE_TOP + Cm(0.3),
        kind_label,
        bg_color=accent,
        text_color=WHITE,
    )

    add_top_line(slide, color=accent)

    # Left accent bar
    add_accent_bar(
        slide,
        MARGIN_LEFT, CONTENT_TOP,
        ACCENT_BAR_WIDTH, Cm(8),
        accent,
    )

    # Info key-value pairs
    info_left = MARGIN_LEFT + Cm(1.2)
    info_width = Cm(28)
    y = CONTENT_TOP + Cm(0.3)
    row_h = Cm(1.6)

    parent_label = ""
    if node["parent"] and node["parent"] in nodes:
        parent_label = nodes[node["parent"]]["label"]
    cond_str = resolve_conditions(node["conditions"], concept_map)
    child_labels = [nodes[c]["label"] for c in node["children"] if c in nodes]
    tag_names = node.get("tag_names", [])

    fields = [
        ("Parent view", parent_label or "— (root)"),
        ("Activated by", cond_str or "— (always)"),
        ("Child views", ", ".join(child_labels) if child_labels else "— (none)"),
    ]
    if tag_names:
        fields.append(("Concepts", ", ".join(tag_names)))

    for label, value in fields:
        # Label
        add_textbox(
            slide,
            info_left, y,
            Cm(6), Cm(0.7),
            label,
            font_size=FONT_SIZE_BODY_SMALL,
            bold=True,
            color=MUTED,
        )
        # Value
        add_textbox(
            slide,
            info_left + Cm(6.5), y,
            info_width - Cm(6.5), Cm(0.7),
            value,
            font_size=FONT_SIZE_BODY,
            color=DARK_TEXT,
        )
        y += row_h

    # Divider line below info
    add_line(
        slide,
        info_left, y,
        MARGIN_LEFT + info_width, y,
        color=DIVIDER,
    )

    # Instruction placeholder text
    y += Cm(0.8)
    add_textbox(
        slide,
        info_left, y,
        info_width, Cm(4),
        "Add annotation instructions for this view here.",
        font_size=FONT_SIZE_BODY,
        italic=True,
        color=MUTED,
    )

    add_bottom_strip(slide, accent)
    return slide


def build_good_examples_slide(prs, node: dict, images_dir=None):
    """Good examples slide — green accent, up to 3 images per slide."""
    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, WHITE)

    accent = kind_color(node["kind"])

    # Title with checkmark
    add_textbox(
        slide,
        TITLE_LEFT, TITLE_TOP,
        TITLE_WIDTH, TITLE_HEIGHT,
        f"✓  Good Examples  —  {node['label']}",
        font_size=FONT_SIZE_SLIDE_TITLE,
        bold=True,
        color=GREEN,
    )

    # Green accent bar under title
    add_accent_bar(
        slide,
        MARGIN_LEFT, TITLE_TOP + TITLE_HEIGHT + Cm(0.1),
        SLIDE_WIDTH - 2 * MARGIN_LEFT, Cm(0.15),
        GREEN,
    )

    # Image gallery area
    gallery_top = CONTENT_TOP + Cm(0.5)
    gallery_height = SLIDE_HEIGHT - gallery_top - Cm(2.0)
    gallery_width = SLIDE_WIDTH - 2 * MARGIN_LEFT

    # Find images for this view
    view_images = find_view_images(node["label"], images_dir)
    tag_names = node.get("tag_names", [])
    concepts = list(tag_names) if tag_names else list(view_images.keys()) if view_images else []

    # If no concepts, show generic placeholders
    if not concepts:
        concepts = ["Example 1", "Example 2", "Example 3"]

    # Max 3 per slide — take first 3
    display_concepts = concepts[:3]
    positions = grid_positions(
        len(display_concepts),
        MARGIN_LEFT, gallery_top,
        gallery_width, gallery_height,
        cols=min(len(display_concepts), 3),
        h_padding=Cm(1.0),
        v_padding=Cm(0.5),
    )

    for (x, y, w, h), concept_name in zip(positions, display_concepts):
        # Reserve space for label below image
        label_h = Cm(1.2)
        img_h = h - label_h

        # Try to find actual image
        img_path = _match_image(concept_name, view_images)

        if img_path and img_path.exists():
            add_image(slide, img_path, x, y, w, img_h)
        else:
            add_image_placeholder(
                slide, x, y, w, img_h,
                label="[ good example ]",
                border_color=GREEN,
                bg_color=PLACEHOLDER_BG,
            )

        # Concept label
        add_textbox(
            slide, x, y + img_h, w, label_h,
            concept_name,
            font_size=FONT_SIZE_BODY_SMALL,
            bold=True,
            color=DARK_TEXT,
            alignment=PP_ALIGN.CENTER,
        )

    add_bottom_strip(slide, GREEN)
    return slide


def build_bad_examples_slide(prs, node: dict, images_dir=None):
    """Bad examples slide — red accent, placeholders only for now."""
    slide = add_blank_slide(prs)
    set_slide_bg_solid(slide, WHITE)

    # Title with X mark
    add_textbox(
        slide,
        TITLE_LEFT, TITLE_TOP,
        TITLE_WIDTH, TITLE_HEIGHT,
        f"✗  Bad Examples  —  {node['label']}",
        font_size=FONT_SIZE_SLIDE_TITLE,
        bold=True,
        color=RED,
    )

    # Red accent bar under title
    add_accent_bar(
        slide,
        MARGIN_LEFT, TITLE_TOP + TITLE_HEIGHT + Cm(0.1),
        SLIDE_WIDTH - 2 * MARGIN_LEFT, Cm(0.15),
        RED,
    )

    # Placeholder gallery
    gallery_top = CONTENT_TOP + Cm(0.5)
    gallery_height = SLIDE_HEIGHT - gallery_top - Cm(2.0)
    gallery_width = SLIDE_WIDTH - 2 * MARGIN_LEFT

    tag_names = node.get("tag_names", [])
    concepts = list(tag_names) if tag_names else ["Example 1", "Example 2", "Example 3"]

    display_concepts = concepts[:3]
    positions = grid_positions(
        len(display_concepts),
        MARGIN_LEFT, gallery_top,
        gallery_width, gallery_height,
        cols=min(len(display_concepts), 3),
        h_padding=Cm(1.0),
        v_padding=Cm(0.5),
    )

    for (x, y, w, h), concept_name in zip(positions, display_concepts):
        label_h = Cm(1.2)
        img_h = h - label_h

        add_image_placeholder(
            slide, x, y, w, img_h,
            label="[ bad example ]",
            border_color=RED,
            bg_color=PLACEHOLDER_BG,
        )

        add_textbox(
            slide, x, y + img_h, w, label_h,
            concept_name,
            font_size=FONT_SIZE_BODY_SMALL,
            bold=True,
            color=DARK_TEXT,
            alignment=PP_ALIGN.CENTER,
        )

    add_bottom_strip(slide, RED)
    return slide


# ──────────────────────────────────────────────────────────────────────
# Image matching helper
# ──────────────────────────────────────────────────────────────────────


def _match_image(concept_name: str, view_images: dict[str, Path]) -> Path | None:
    """Try to match a concept name to a downloaded image."""
    if not view_images:
        return None

    # Direct match
    if concept_name in view_images:
        return view_images[concept_name]

    # Sanitized match
    sanitized = _sanitize(concept_name).replace("_", " ")
    for k, v in view_images.items():
        if k == sanitized or k.lower() == concept_name.lower():
            return v

    # Partial match
    cn_lower = concept_name.lower().replace(" ", "_")
    for k, v in view_images.items():
        if cn_lower in k.lower().replace(" ", "_"):
            return v

    return None


# ══════════════════════════════════════════════════════════════════════
# ORCHESTRATOR — builds the full deck
# ══════════════════════════════════════════════════════════════════════


def build_all_slides(
    prs,
    project_map: dict,
    *,
    images_dir=None,
    org_slug: str = "",
    project_slug: str = "",
):
    """Build the complete annotation guide slide deck.

    This is the main entry point. Edit this function to change
    which slides appear and in what order.
    """
    nodes, roots = build_tree(project_map)
    concept_map = build_concept_map(project_map)
    ordered = dfs_order(nodes, roots)

    # ── Cover ──
    project_name = " — ".join(filter(None, [org_slug.upper(), project_slug.upper()]))
    build_cover_slide(prs, project_name=project_name)

    # ── Table of Contents ──
    build_toc_slide(prs, nodes, roots)

    # ── Views Overview (tree diagram) ──
    build_overview_slide(prs, nodes, roots, concept_map)

    # ── Per-view slides ──
    for nid in ordered:
        node = nodes[nid]
        is_root = not node["parent"] or node["parent"] not in nodes

        # Section divider for root views only
        if is_root:
            build_section_slide(prs, node["label"], node["kind"])

        # Info slide
        build_info_slide(prs, node, nodes, concept_map)

        # Good examples
        build_good_examples_slide(prs, node, images_dir=images_dir)

        # Bad examples
        build_bad_examples_slide(prs, node, images_dir=images_dir)

    logger.info("Built %d slides total.", len(prs.slides))
